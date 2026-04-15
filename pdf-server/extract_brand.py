#!/usr/bin/env python3
"""
extract_brand.py — Extract brand assets from a PPTX file.

Extracts: colors (theme + frequency), fonts, logo candidates, header/footer,
dark/light mode detection. Outputs JSON to stdout.

Usage: python3 extract_brand.py <input.pptx> [output_dir]

Dependencies: pip install python-pptx Pillow
"""

import sys
import os
import json
import base64
import hashlib
from collections import Counter
from io import BytesIO

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# ─── Helpers ───────────────────────────────────────────────────────────────────

def rgb_to_hex(r, g, b):
    return f"#{r:02x}{g:02x}{b:02x}"

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

def luminance(hex_color):
    """Relative luminance (0=black, 1=white)"""
    r, g, b = hex_to_rgb(hex_color)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255

def is_near_black(hex_color, threshold=0.15):
    return luminance(hex_color) < threshold

def is_near_white(hex_color, threshold=0.90):
    return luminance(hex_color) > threshold

# ─── Theme Color Extraction ───────────────────────────────────────────────────

def extract_theme_colors(prs):
    """Extract accent1-6, dk1, dk2, lt1, lt2 from the slide master theme."""
    theme_colors = {}
    try:
        master = prs.slide_masters[0]
        theme_xml = master.element.find(".//" + qn("a:clrScheme"))
        if theme_xml is not None:
            for child in theme_xml:
                tag_name = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                # Get the color value
                srgb = child.find(qn("a:srgbClr"))
                sys_clr = child.find(qn("a:sysClr"))
                if srgb is not None:
                    theme_colors[tag_name] = f"#{srgb.get('val', '000000')}"
                elif sys_clr is not None:
                    theme_colors[tag_name] = f"#{sys_clr.get('lastClr', '000000')}"
    except Exception:
        pass
    return theme_colors

# ─── Color Frequency Scan ─────────────────────────────────────────────────────

def extract_color_frequency(prs):
    """Scan all shapes for fill colors and text colors, count frequency."""
    color_counter = Counter()

    for slide in prs.slides:
        for shape in slide.shapes:
            # Shape fill color
            try:
                if shape.fill and shape.fill.type is not None:
                    fc = shape.fill.fore_color
                    if fc and fc.rgb:
                        hex_color = f"#{fc.rgb}"
                        color_counter[hex_color.lower()] += 1
            except Exception:
                pass

            # Text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                hex_color = f"#{run.font.color.rgb}"
                                color_counter[hex_color.lower()] += 1
                        except Exception:
                            pass

    return color_counter

# ─── Font Extraction ──────────────────────────────────────────────────────────

def extract_fonts(prs):
    """Extract font families and their usage frequency."""
    font_counter = Counter()
    font_sizes = {}  # font_name → list of sizes

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        name = run.font.name
                        if name:
                            font_counter[name] += 1
                            size = run.font.size
                            if size:
                                pt = size.pt
                                if name not in font_sizes:
                                    font_sizes[name] = []
                                font_sizes[name].append(pt)
                    except Exception:
                        pass

    # Determine primary/secondary
    ranked = font_counter.most_common()
    primary = ranked[0][0] if ranked else None
    secondary = ranked[1][0] if len(ranked) > 1 else None

    by_frequency = [{"name": name, "count": count} for name, count in ranked[:5]]

    return {
        "primary": primary,
        "secondary": secondary if secondary != primary else None,
        "by_frequency": by_frequency,
    }

# ─── Logo Extraction ──────────────────────────────────────────────────────────

def extract_logos(prs, output_dir=None):
    """Find logo candidates — small images in the top/corner of early slides."""
    candidates = []
    seen_hashes = set()

    # Check slide masters first (logos are almost always here), then first 3 slides
    check_sources = []
    for master in prs.slide_masters:
        check_sources.append((-1, master))  # -1 = slide master
    for i, slide in enumerate(list(prs.slides)[:3]):
        check_sources.append((i, slide))

    for slide_idx, slide in check_sources:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            try:
                img = shape.image
                blob = img.blob
                img_hash = hashlib.md5(blob).hexdigest()

                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)

                # Heuristic: logo is small and positioned in top or corner
                w_inches = shape.width / 914400
                h_inches = shape.height / 914400
                left_inches = shape.left / 914400
                top_inches = shape.top / 914400

                # Logo heuristic: small image (< 3 inches) near top (< 2 inches from top)
                is_small = w_inches < 3 and h_inches < 2
                is_top = top_inches < 2
                is_corner = left_inches < 2 or left_inches > 8  # left or right corner

                if not (is_small and is_top):
                    continue

                confidence = "high" if (is_corner and is_small) else "medium"

                # Save image if output dir provided
                filename = f"logo_{len(candidates)}.{img.content_type.split('/')[-1]}"
                data_uri = f"data:{img.content_type};base64,{base64.b64encode(blob).decode()}"

                if output_dir:
                    os.makedirs(os.path.join(output_dir, "images"), exist_ok=True)
                    with open(os.path.join(output_dir, "images", filename), "wb") as f:
                        f.write(blob)

                candidates.append({
                    "filename": filename,
                    "data_uri": data_uri,
                    "confidence": confidence,
                    "size": {"w": round(w_inches, 2), "h": round(h_inches, 2)},
                    "position": {"x": round(left_inches, 2), "y": round(top_inches, 2)},
                    "slide_index": slide_idx,
                })

            except Exception:
                pass

    return candidates

# ─── Header/Footer Extraction ─────────────────────────────────────────────────

def extract_header_footer(prs):
    """Detect header/footer elements: legal text, slide numbers, dates, dividers."""
    header = {"has_content": False, "elements": []}
    footer = {
        "has_content": False,
        "elements": [],
        "builtin": {
            "footer_placeholder": None,
            "slide_number_placeholder": None,
            "date_placeholder": None,
        },
    }

    slide_height_inches = prs.slide_height / 914400

    for slide in prs.slides:
        for shape in slide.placeholders:
            try:
                idx = shape.placeholder_format.idx
                # Footer placeholder (idx 11)
                if idx == 11 and shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        footer["has_content"] = True
                        footer["builtin"]["footer_placeholder"] = {"text": text}
                # Slide number placeholder (idx 12)
                elif idx == 12:
                    footer["has_content"] = True
                    footer["builtin"]["slide_number_placeholder"] = {"text": ""}
                # Date placeholder (idx 10)
                elif idx == 10:
                    footer["builtin"]["date_placeholder"] = {"text": ""}
            except Exception:
                pass

        # Scan for text in top 15% zone (header) and bottom 15% zone (footer)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                top_inches = shape.top / 914400
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Top zone = header
                if top_inches < slide_height_inches * 0.15:
                    header["has_content"] = True
                    if text not in [e.get("text") for e in header["elements"]]:
                        header["elements"].append({"type": "text", "text": text})

                # Bottom zone = footer
                if top_inches > slide_height_inches * 0.80:
                    # Check for legal patterns
                    lower = text.lower()
                    if any(kw in lower for kw in ["confidential", "©", "proprietary", "copyright", "all rights"]):
                        footer["has_content"] = True
                        footer["elements"].append({"type": "legal_text", "text": text})
                    elif not footer["builtin"]["footer_placeholder"]:
                        footer["has_content"] = True
                        footer["elements"].append({"type": "text", "text": text})
            except Exception:
                pass

    # Detect divider lines (thin + wide shapes)
    divider_lines = []
    for slide in list(prs.slides)[:3]:
        for shape in slide.shapes:
            try:
                w_inches = shape.width / 914400
                h_inches = shape.height / 914400
                top_inches = shape.top / 914400

                # Divider: very thin height (< 0.1") and wide (> 5")
                if h_inches < 0.1 and w_inches > 5:
                    zone = "header" if top_inches < slide_height_inches * 0.2 else "footer"
                    divider_lines.append({
                        "zone": zone,
                        "y_inches": round(top_inches, 2),
                    })
            except Exception:
                pass

    return {
        "header": header,
        "footer": footer,
        "divider_lines": divider_lines,
    }

# ─── Background / Dark Mode Detection ─────────────────────────────────────────

def detect_mode(prs):
    """Detect if the majority of slides use dark or light backgrounds."""
    dark_count = 0
    light_count = 0

    for slide in prs.slides:
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                fc = fill.fore_color
                if fc and fc.rgb:
                    hex_color = f"#{fc.rgb}"
                    if is_near_black(hex_color):
                        dark_count += 1
                    elif is_near_white(hex_color):
                        light_count += 1
                    else:
                        light_count += 1  # non-extreme colors count as light
        except Exception:
            light_count += 1  # default to light if can't detect

    return "dark" if dark_count > light_count else "light"

# ─── Main ──────────────────────────────────────────────────────────────────────

def extract_brand(pptx_path, output_dir=None):
    prs = Presentation(pptx_path)

    # Theme colors
    theme_colors = extract_theme_colors(prs)

    # Color frequency
    color_counter = extract_color_frequency(prs)

    # Filter out near-black and near-white for primary/secondary detection
    meaningful_colors = {
        c: n for c, n in color_counter.items()
        if not is_near_black(c) and not is_near_white(c)
    }
    ranked_colors = sorted(meaningful_colors.items(), key=lambda x: -x[1])

    primary_color = ranked_colors[0][0] if ranked_colors else None
    secondary_color = ranked_colors[1][0] if len(ranked_colors) > 1 else None

    by_frequency = [
        {"hex": color, "count": count}
        for color, count in color_counter.most_common(10)
    ]

    # Fonts
    fonts = extract_fonts(prs)

    # Logos
    logo_candidates = extract_logos(prs, output_dir)

    # Header/Footer
    header_footer = extract_header_footer(prs)

    # Dark/Light mode
    mode = detect_mode(prs)

    # Build brand.json
    brand = {
        "colors": {
            "primary": primary_color,
            "secondary": secondary_color,
            "theme": theme_colors,
            "by_frequency": by_frequency,
            "mode": mode,
        },
        "fonts": fonts,
        "logo_candidates": logo_candidates,
        "header_footer": header_footer,
        "extraction_status": "complete",
    }

    # Save if output dir provided
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
        with open(os.path.join(output_dir, "brand.json"), "w") as f:
            json.dump(brand, f, indent=2)

    return brand


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 extract_brand.py <input.pptx> [output_dir]", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None

    brand = extract_brand(pptx_path, output_dir)
    print(json.dumps(brand))
